import os
import requests
from pptx import Presentation
from utils import remove_all_slides
from logger import LOG
from urllib.parse import urlparse
import hashlib
from typing import Optional

def download_image(image_url: str, save_dir: str = "temp_images") -> Optional[str]:
    """
    下載圖片並返回本地保存路徑
    
    Args:
        image_url (str): 圖片URL
        save_dir (str): 保存目錄
        
    Returns:
        Optional[str]: 本地文件路徑，如果下載失敗則返回 None
    """
    try:
        # 創建保存目錄
        os.makedirs(save_dir, exist_ok=True)
        
        # 使用 URL 生成唯一的文件名
        url_hash = hashlib.md5(image_url.encode()).hexdigest()
        # 從 URL 中獲取文件擴展名
        parsed_url = urlparse(image_url)
        file_ext = os.path.splitext(parsed_url.path)[1]
        if not file_ext:
            file_ext = '.jpg'  # 默認擴展名
        
        # 構建本地文件路徑
        local_path = os.path.join(save_dir, f"{url_hash}{file_ext}")
        
        # 如果文件已存在，直接返回路徑
        if os.path.exists(local_path):
            return local_path
            
        # 下載圖片
        response = requests.get(image_url, timeout=10)
        response.raise_for_status()
        
        # 保存圖片
        with open(local_path, 'wb') as f:
            f.write(response.content)
            
        LOG.info(f"圖片已下載到: {local_path}")
        return local_path
        
    except Exception as e:
        LOG.error(f"下載圖片時發生錯誤: {str(e)}")
        return None

def generate_presentation(powerpoint_data, template_path: str, output_path: str):
    """生成 PowerPoint 演示文稿"""
    # 檢查模板文件是否存在
    if not os.path.exists(template_path):
        LOG.error(f"模板文件 '{template_path}' 不存在。")
        raise FileNotFoundError(f"模板文件 '{template_path}' 不存在。")

    prs = Presentation(template_path)
    remove_all_slides(prs)
    prs.core_properties.title = powerpoint_data.title

    # 遍歷所有幻燈片數據
    for slide in powerpoint_data.slides:
        # 確保布局索引不超出範圍
        if slide.layout_id >= len(prs.slide_layouts):
            slide_layout = prs.slide_layouts[0]
        else:
            slide_layout = prs.slide_layouts[slide.layout_id]

        new_slide = prs.slides.add_slide(slide_layout)

        # 設置幻燈片標題
        if new_slide.shapes.title:
            new_slide.shapes.title.text = slide.content.title
            LOG.debug(f"設置幻燈片標題: {slide.content.title}")

        # 添加文本內容
        for shape in new_slide.shapes:
            if shape.has_text_frame and not shape == new_slide.shapes.title:
                text_frame = shape.text_frame
                text_frame.clear()
                for point in slide.content.bullet_points:
                    p = text_frame.add_paragraph()
                    p.text = point
                    p.level = 0
                    LOG.debug(f"添加列表項: {point}")
                break

        # 處理圖片
        if slide.content.image_path:
            # 檢查是否是 URL
            if slide.content.image_path.startswith(('http://', 'https://')):
                # 下載圖片
                local_image_path = download_image(slide.content.image_path)
                if local_image_path:
                    image_path = local_image_path
                else:
                    LOG.warning(f"無法下載圖片: {slide.content.image_path}")
                    continue
            else:
                # 如果不是 URL，假設是本地路徑
                image_path = os.path.join(os.getcwd(), slide.content.image_path)

            # 檢查圖片文件是否存在
            if os.path.exists(image_path):
                # 插入圖片到占位符中
                for shape in new_slide.placeholders:
                    if shape.placeholder_format.type == 18:  # 圖片占位符
                        try:
                            shape.insert_picture(image_path)
                            LOG.debug(f"插入圖片: {image_path}")
                            break
                        except Exception as e:
                            LOG.error(f"插入圖片時發生錯誤: {str(e)}")
            else:
                LOG.warning(f"圖片路徑不存在: {image_path}")

    # 保存生成的 PowerPoint 文件
    try:
        prs.save(output_path)
        LOG.info(f"演示文稿已保存到: {output_path}")
    except Exception as e:
        LOG.error(f"保存演示文稿時發生錯誤: {str(e)}")
        raise

